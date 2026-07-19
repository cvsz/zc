"""
zc_powerpoint.py — Conversational slide-deck assistant
AI Model Coder CLI v1.15.0

A chat loop purpose-built for slide decks: describe what you want in plain
English — add a slide, restyle a title, turn bullets into a comparison
table, drop in a chart — and each turn updates an in-memory list of slides
that gets written out to a real .pptx file after every turn, so the deck
on disk always reflects the latest state of the conversation.

This is the CLI's answer to a "zAICoder in PowerPoint"-style experience: no
Office add-in, just a terminal chat that keeps a real .pptx file in sync.
Mirrors zc_excel.py's design one-for-one (same session/history/undo
shape, same denylisted-exec safety model) since both are the same kind of
product: a chat loop that edits a live in-memory document and re-saves the
real file after every turn.

Requires: python-pptx (see requirements.txt — optional dep, only needed
for this module).

CLI:
  --pptx [FILE]            Start a PowerPoint chat session, optionally
                            loading an existing .pptx as the starting deck
  --pptx-output FILE       Deck to write after every turn
                            (default: <input>.pptx or pptx_session.pptx)

Slash commands inside the session:
  /help                Show this help
  /exit, /quit         End the session (deck is already saved)
  /slides              List current slides and their titles
  /show N              Print the text content of slide N
  /undo                Revert to the state before the last applied change
"""
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None  # type: ignore[assignment]

SYSTEM_PROMPT = """\
You are a slide-deck assistant embedded in a CLI tool. The user is \
chatting with you to add slides, restyle text, turn bullets into tables, \
and add charts — all applied directly to a live in-memory deck \
representation that gets written back to a real .pptx file after every \
turn.

You have access to a list called `slides`, where each slide is a dict:
  {"title": str, "bullets": [str, ...], "layout": "title_content" | "title_only" | "section_header",
   "table": {"headers": [str,...], "rows": [[str,...], ...]} | None,
   "chart": {"type": "bar"|"line"|"pie", "categories": [str,...], "series": {name: [num,...]}} | None}

To change the deck, call the provided helper functions — do not build the
.pptx file yourself:
  add_slide(title, bullets=None, layout="title_content", table=None, chart=None)
  update_slide(index, title=None, bullets=None, table=None, chart=None)
  delete_slide(index)
  reorder_slides(new_order)   # new_order is a list of old indices, e.g. [0,2,1]

Respond in ONE of two ways:

1. If the request requires changing the deck (adding/editing/removing/
reordering slides), respond with ONLY a single fenced python code block
that calls the helper functions above. No prose outside the code block.

2. If the request is a question that doesn't require changing the deck
(e.g. "how many slides do I have", "what does slide 3 say"), answer in
plain text with no code block, using the deck summary given to you.

Never write to disk, never import anything, never use `open(`, `os`,
`sys`, `subprocess`, or `eval`/`exec`. Keep bullets concise (under ~12
words each) and cap slides at 5-6 bullets — this is a presentation, not a
document.
"""

_CODE_BLOCK = re.compile(r"```(?:python)?\s*\n(.*?)```", re.DOTALL)

# Mirrors zc_excel.py's denylist for the same reason: this is a
# best-effort catch for obviously unsafe generated code, not a real
# sandbox. Anything more sensitive should go through
# --code-agent-sandbox instead.
_DENYLIST = (
    "import os", "import sys", "import subprocess", "import socket",
    "__import__", "open(", "eval(", "exec(", "os.", "sys.", "subprocess.",
    "socket.", "shutil.", ".system(", "pathlib",
)


class PptxSession:
    def __init__(self, input_path=None):
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for --pptx (pip install python-pptx)"
            )
        self.slides = []
        self._history_stack = []  # for /undo — list of deep-copied slide lists
        self._template_path = None  # if loaded from an existing deck, reuse its theme

        if input_path:
            self._load(input_path)

    def _load(self, path):
        prs = Presentation(path)
        self._template_path = path
        for slide in prs.slides:
            title = ""
            bullets = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if not title and shape == slide.shapes.title:
                    title = text
                elif text:
                    bullets.extend(line for line in text.split("\n") if line.strip())
            self.slides.append({
                "title": title, "bullets": bullets, "layout": "title_content",
                "table": None, "chart": None,
            })

    # ── context for the model ───────────────────────────────────────────

    def summary(self):
        if not self.slides:
            return "(no slides yet)"
        parts = [f"Deck has {len(self.slides)} slide(s):"]
        for i, s in enumerate(self.slides):
            extra = []
            if s.get("table"):
                extra.append(f"table {len(s['table']['rows'])}x{len(s['table']['headers'])}")
            if s.get("chart"):
                extra.append(f"{s['chart']['type']} chart")
            extra_str = f" [{', '.join(extra)}]" if extra else ""
            bullets_preview = "; ".join(s["bullets"][:3])
            parts.append(f"  {i}: \"{s['title']}\"{extra_str} — {bullets_preview}")
        return "\n".join(parts)

    # ── applying a model turn ───────────────────────────────────────────

    def _snapshot(self):
        import copy
        self._history_stack.append(copy.deepcopy(self.slides))
        if len(self._history_stack) > 20:
            self._history_stack.pop(0)

    def undo(self):
        if not self._history_stack:
            return False
        self.slides = self._history_stack.pop()
        return True

    def apply_code(self, code):
        """Run model-generated code against `self.slides`. Returns (ok, message)."""
        lowered = code.lower()
        for bad in _DENYLIST:
            if bad in lowered:
                return False, f"[blocked] generated code used a disallowed construct: {bad!r}"

        self._snapshot()

        def add_slide(title, bullets=None, layout="title_content", table=None, chart=None):
            self.slides.append({
                "title": title, "bullets": bullets or [], "layout": layout,
                "table": table, "chart": chart,
            })

        def update_slide(index, title=None, bullets=None, table=None, chart=None):
            s = self.slides[index]
            if title is not None:
                s["title"] = title
            if bullets is not None:
                s["bullets"] = bullets
            if table is not None:
                s["table"] = table
            if chart is not None:
                s["chart"] = chart

        def delete_slide(index):
            self.slides.pop(index)

        def reorder_slides(new_order):
            self.slides = [self.slides[i] for i in new_order]

        local_ns = {
            "slides": self.slides,
            "add_slide": add_slide,
            "update_slide": update_slide,
            "delete_slide": delete_slide,
            "reorder_slides": reorder_slides,
        }
        try:
            exec(compile(code, "<pptx-turn>", "exec"), {"__builtins__": {
                "len": len, "range": range, "sum": sum, "min": min, "max": max,
                "round": round, "sorted": sorted, "list": list, "dict": dict,
                "str": str, "int": int, "float": float, "bool": bool,
                "enumerate": enumerate, "zip": zip, "abs": abs,
            }}, local_ns)
        except Exception as e:
            self.undo()
            return False, f"[ERROR] generated code failed: {e}"
        return True, "applied"

    # ── persistence ──────────────────────────────────────────────────────

    def save(self, output_path):
        # Rebuilt from scratch on every save (simplest correct approach given
        # `slides` is the single source of truth) rather than diffing against
        # a live Presentation object — mirrors zc_excel.py rewriting the
        # whole workbook from `sheets` on every save.
        prs = Presentation()
        title_content = prs.slide_layouts[1]
        title_only = prs.slide_layouts[5]
        section_header = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else title_only

        for s in self.slides:
            layout = {"title_only": title_only, "section_header": section_header}.get(
                s.get("layout"), title_content)
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = s["title"]

            if s.get("bullets") and layout is title_content:
                body = None
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        body = ph
                        break
                if body is not None:
                    tf = body.text_frame
                    tf.clear()
                    for i, bullet in enumerate(s["bullets"]):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = bullet

            if s.get("table"):
                self._add_table(slide, s["table"])
            if s.get("chart"):
                self._add_chart(slide, s["chart"])

        prs.save(output_path)

    def _add_table(self, slide, table):
        headers, rows = table["headers"], table["rows"]
        n_rows, n_cols = len(rows) + 1, len(headers)
        left, top, width, height = Inches(0.5), Inches(1.8), Inches(9), Inches(0.4 * n_rows)
        shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        tbl = shape.table
        for c, header in enumerate(headers):
            tbl.cell(0, c).text = str(header)
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                tbl.cell(r, c).text = str(val)

    def _add_chart(self, slide, chart):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        xl_type = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }.get(chart["type"], XL_CHART_TYPE.COLUMN_CLUSTERED)

        data = CategoryChartData()
        data.categories = chart["categories"]
        for name, values in chart["series"].items():
            data.add_series(name, values)

        left, top, width, height = Inches(1), Inches(1.8), Inches(8), Inches(4.5)
        slide.shapes.add_chart(xl_type, left, top, width, height, data)


HELP_TEXT = """\
Commands:
  /help                Show this help
  /exit, /quit         End the session (deck is already saved)
  /slides              List current slides and their titles
  /show N              Print the text content of slide N
  /undo                Revert to the state before the last applied change
"""


def cmd_pptx_chat(api_key, model, input_path=None, output_path=None,
                   temperature=0.3, max_tokens=4096, native=False):
    """native=True routes each turn through zc_skills_api.py's pptx
    Skill (Anthropic's own maintained implementation, server-side in a
    code-execution container) instead of the hand-rolled python-pptx path
    below. See --pptx-native. Requires Skills access on the account; the
    hand-rolled path here remains the default and the fallback for
    accounts without it."""
    if native:
        return _cmd_pptx_chat_native(api_key, model, input_path=input_path,
                                     output_path=output_path, max_tokens=max_tokens)

    if Presentation is None:
        print("[ERROR] python-pptx is required for --pptx. Install with: "
              "pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    from wire.coder import Coder

    try:
        session = PptxSession(input_path)
    except (ImportError, ValueError, OSError) as e:
        print(f"[ERROR] {e}", file=sys.stderr)
        sys.exit(1)

    output_path = output_path or (
        re.sub(r"\.\w+$", "", input_path) + ".pptx" if input_path else "pptx_session.pptx"
    )

    c = Coder(api_key=api_key, model=model, temperature=temperature, max_tokens=max_tokens)

    print(f"\033[94mAI Model Coder — PowerPoint chat\033[0m  (model: {c.model})")
    print(f"Deck: {output_path}  (saved after every applied change)")
    print("Type /help for commands, /exit to quit.\n")

    history = []
    while True:
        try:
            user_input = input("\033[92myou›\033[0m ").strip()
        except (EOFError, KeyboardInterrupt):
            print()
            break
        if not user_input:
            continue

        if user_input.startswith("/"):
            parts = user_input.split()
            cmd = parts[0].lower()
            if cmd in ("/exit", "/quit"):
                break
            if cmd == "/help":
                print(HELP_TEXT); continue
            if cmd == "/slides":
                for i, s in enumerate(session.slides):
                    print(f"  {i}: {s['title']!r} ({len(s['bullets'])} bullets)")
                continue
            if cmd == "/show":
                if len(parts) < 2 or not parts[1].isdigit() or int(parts[1]) >= len(session.slides):
                    print(f"[usage] /show N — have {len(session.slides)} slide(s)")
                    continue
                s = session.slides[int(parts[1])]
                print(f"Title: {s['title']}")
                for b in s["bullets"]:
                    print(f"  - {b}")
                continue
            if cmd == "/undo":
                print("[reverted]" if session.undo() else "[nothing to undo]")
                session.save(output_path)
                continue
            print(f"[unknown command {cmd!r}; try /help]")
            continue

        prompt = f"Current deck:\n{session.summary()}\n\nRequest: {user_input}"
        reply = c.generate(prompt, system=SYSTEM_PROMPT, history=history)

        match = _CODE_BLOCK.search(reply)
        if match:
            ok, message = session.apply_code(match.group(1))
            if ok:
                session.save(output_path)
                print(f"\033[96mzc›\033[0m Updated and saved to {output_path} "
                     f"({len(session.slides)} slides)\n")
            else:
                print(f"\033[93mzc›\033[0m {message}\n")
        else:
            print(f"\033[96mzc›\033[0m {reply}\n")

        history.append({"role": "user", "content": user_input})
        history.append({"role": "assistant", "content": reply})

    print(f"\033[94mSession ended. Final deck: {output_path}\033[0m")
    return output_path


def _cmd_pptx_chat_native(api_key, model, input_path=None, output_path=None, max_tokens=4096):
    """--pptx-native: mirrors zc_excel.py's _cmd_excel_chat_native
    one-for-one (same reasoning applies here — see that function's
    docstring), just against the pptx Skill instead of xlsx. No local
    python-pptx dependency needed for this path.

    Slash commands from the hand-rolled path (/slides, /show, /undo)
    aren't available here — the pptx Skill owns the deck, this CLI has no
    local copy of it to inspect or revert.
    """
    from wire.zc_files import FilesAPI
    from wire.zc_skills_api import SkillsApiClient, build_user_content, extract_output_file_ids

    files_api = FilesAPI(api_key=api_key, model=model)
    client = SkillsApiClient(api_key=api_key, model=model, max_tokens=max_tokens)

    output_path = output_path or (
        re.sub(r"\.\w+$", "", input_path) + ".pptx" if input_path else "pptx_session.pptx"
    )

    pending_file_ids = []
    if input_path:
        try:
            uploaded = files_api.upload(input_path)
        except (RuntimeError, OSError) as e:
            print(f"[ERROR] Could not upload {input_path}: {e}", file=sys.stderr)
            sys.exit(1)
        fid = uploaded.get("id")
        if not fid:
            print(f"[ERROR] Upload succeeded but returned no file id: {uploaded}", file=sys.stderr)
            sys.exit(1)
        pending_file_ids = [fid]

    print(f"\033[94mAI Model Coder — PowerPoint chat (native Skills API)\033[0m  (model: {model})")
    print(f"Deck: {output_path}  (saved after every turn that produces one)")
    print("Type /exit to quit. (/slides, /show, /undo aren't available in --pptx-native.)\n")

    messages = []
    container_id = None
    while True:
        try:
            user_input = input("\033[92myou›\033[0m ").strip()
        except (EOFError, KeyboardInterrupt):
            print()
            break
        if not user_input:
            continue
        if user_input.lower() in ("/exit", "/quit"):
            break

        messages.append({"role": "user", "content": build_user_content(user_input, pending_file_ids)})
        has_uploads = bool(pending_file_ids)
        pending_file_ids = []  # only attach on the turn that actually introduces the file

        data = client.call_with_skills_turn(
            messages, skills=["pptx"], container_id=container_id, has_file_uploads=has_uploads,
        )
        if "error" in data:
            print(f"\033[91m✗ {data['error']}\033[0m\n")
            messages.pop()
            continue

        container_id = (data.get("container") or {}).get("id", container_id)
        messages.append({"role": "assistant", "content": data.get("content", [])})

        text = "".join(b.get("text", "") for b in data.get("content", []) if b.get("type") == "text")
        if text:
            print(f"\033[96mzc›\033[0m {text}\n")

        new_file_ids = extract_output_file_ids(data)
        if new_file_ids:
            try:
                files_api.download(new_file_ids[-1], output_path)
                print(f"\033[90m  (saved to {output_path})\033[0m\n")
            except RuntimeError as e:
                print(f"\033[93m  Couldn't download generated file: {e}\033[0m\n")

    print(f"\033[94mSession ended. Final deck: {output_path}\033[0m")
    return output_path